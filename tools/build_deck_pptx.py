#!/usr/bin/env python3
"""Build the RVC legislator deck as .pptx in the Editorial DS (jeffpinto.com design system).

Palette/typography per 'The Editorial DS — portable design system' with the policy accent
(#1E3A8A). Fonts mapped to universally-available faces so the deck renders identically in
PowerPoint, Keynote, and Google Slides: Georgia (serif), Arial (sans), Courier New (mono).
Output: site/RVC_Legislator_Deck.pptx — 17 live slides + lean appendix (validated numbers
and sources only; the corrections history lives at /validation.html, not in the deck).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---- Editorial DS tokens -------------------------------------------------
BG     = RGBColor(0xF1, 0xEC, 0xE2)
PAPER  = RGBColor(0xFA, 0xF6, 0xEC)
INK    = RGBColor(0x0E, 0x11, 0x16)
BODY   = RGBColor(0x2A, 0x2A, 0x2A)
MUTED  = RGBColor(0x6B, 0x6B, 0x6B)
RULE   = RGBColor(0x23, 0x1F, 0x1A)
RULESOFT = RGBColor(0xC9, 0xC3, 0xB4)
ACCENT = RGBColor(0x1E, 0x3A, 0x8A)
ACCENT_LT = RGBColor(0x9D, 0xB4, 0xE8)
CREAM_DIM = RGBColor(0xD8, 0xD3, 0xC6)
MONO_DIM  = RGBColor(0xA9, 0xA3, 0x94)
C_STATE = RGBColor(0xA0, 0x7A, 0x2F); T_STATE = RGBColor(0xEC, 0xE2, 0xC9)
C_TAX   = RGBColor(0x4E, 0x7A, 0x58); T_TAX   = RGBColor(0xDE, 0xE7, 0xDC)
C_DIST  = RGBColor(0xB0, 0x54, 0x42); T_DIST  = RGBColor(0xEF, 0xDC, 0xD4)

SERIF, SANS, MONO = "Georgia", "Arial", "Courier New"
W, H = Inches(13.333), Inches(7.5)
MX = Inches(0.62)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
PAGE_N = 0  # auto page number

def slide(dark=False):
    global PAGE_N
    PAGE_N += 1
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = INK if dark else BG; r.line.fill.background()
    r.shadow.inherit = False
    return s

def pno(): return f"{PAGE_N:02d}"

def rect(s, x, y, w, h, color, line=None):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if color is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = color
    if line is None: r.line.fill.background()
    else: r.line.color.rgb = line; r.line.width = Pt(0.75)
    r.shadow.inherit = False
    return r

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing: p.line_spacing = line_spacing
        for (t, fname, size, color, bold, ital, track) in para:
            r = p.add_run(); r.text = t
            f = r.font; f.name = fname; f.size = Pt(size); f.color.rgb = color
            f.bold = bold; f.italic = ital
            if track: r._r.rPr.set('spc', str(int(track * 100)))
    return tb

def R(t, f=SANS, s=12, c=BODY, b=False, i=False, tr=0): return (t, f, s, c, b, i, tr)
def mono(t, s=9, c=MUTED, b=False, tr=1.6): return R(t.upper(), MONO, s, c, b, False, tr)
def eyebrow(t, s=10): return [R("■ ", MONO, s-1, ACCENT, True), mono(t, s, ACCENT, True)]

def masthead(s, left):
    text(s, MX, Inches(0.30), Inches(8), Inches(0.32), [[R(left, SERIF, 15, INK, True)]])
    text(s, W - MX - Inches(5.6), Inches(0.34), Inches(5.6), Inches(0.3),
         [[mono(f"rvc-taxes.jeffpinto.com   ·   {pno()}", 8.5)]], align=PP_ALIGN.RIGHT)
    rect(s, MX, Inches(0.66), W - MX*2, Pt(2.0), RULE)

def footrule(s, note):
    rect(s, MX, Inches(7.02), W - MX*2, Pt(0.9), RULESOFT)
    text(s, MX, Inches(7.10), Inches(9.6), Inches(0.3), [[mono(note, 7.5)]])
    text(s, W - MX - Inches(2.0), Inches(7.10), Inches(2.0), Inches(0.3), [[mono(pno(), 7.5)]], align=PP_ALIGN.RIGHT)

def kpi_cell(s, x, y, w, num_eyebrow, kpi, unit, claim, src, kpi_size=30, tag=None, tagc=None, tagt=None):
    rect(s, x, y, w, Pt(2), RULE)
    yy = y + Inches(0.10)
    if tag:
        tw = Inches(1.85)
        rect(s, x + w - tw, yy, tw, Inches(0.22), tagt, line=tagc)
        text(s, x + w - tw, yy + Inches(0.015), tw, Inches(0.2), [[mono(tag, 7, tagc, True)]], align=PP_ALIGN.CENTER)
    text(s, x, yy, w - Inches(0.1), Inches(0.22), [eyebrow(num_eyebrow, 8)])
    kp = [R(kpi, SERIF, kpi_size, INK, True)]
    if unit: kp.append(R(" " + unit, SERIF, int(kpi_size*0.52), ACCENT, False, True))
    text(s, x, yy + Inches(0.26), w - Inches(0.05), Inches(0.62), [kp])
    text(s, x, yy + Inches(0.26) + Inches(0.58), w - Inches(0.12), Inches(0.78),
         [[R(claim, SANS, 10, BODY)]], line_spacing=1.12)
    text(s, x, yy + Inches(1.62), w - Inches(0.1), Inches(0.22), [[mono(src, 7)]])

def dark_band(s, x, y, w, h):
    return rect(s, x, y, w, h, INK)

# =========================================================================
# 01 — TITLE
s = slide()
text(s, MX, Inches(0.42), Inches(9), Inches(0.3), [[R("Rockville Centre · Housing × Schools", SERIF, 15, INK, True)]])
text(s, W - MX - Inches(4), Inches(0.46), Inches(4), Inches(0.3), [[mono("Briefing · June 2026", 8.5)]], align=PP_ALIGN.RIGHT)
rect(s, MX, Inches(0.80), W - MX*2, Pt(2.2), RULE)
text(s, MX, Inches(1.18), Inches(10), Inches(0.3), [eyebrow("For federal, state & county legislators")])
text(s, MX, Inches(1.55), Inches(11.5), Inches(2.2),
     [[R("Unlock the homes,", SERIF, 58, INK, True)], [R("keep the teachers.", SERIF, 58, INK, True)]], line_spacing=0.98)
text(s, MX, Inches(3.65), Inches(11.0), Inches(1.0),
     [[R("This spring Rockville Centre balanced its school budget by cutting ≈22 teaching positions and 40 teaching assistants — in a village whose median household earns $151,938 and whose homes average $819,000.", SERIF, 16, BODY, False, True)]],
     line_spacing=1.3)
text(s, MX, Inches(4.78), Inches(11.4), Inches(1.5),
     [[R("The problem isn't wealth. It's that the homes stopped changing hands: 27.9% of owner households are headed by someone 65+, enrollment is down ~7% in a decade (≈260 students), and the seniors who want to move face federal tax penalties, benefit-timing traps, and missing downsize inventory. ", SANS, 12.5, BODY),
       R("Everything proposed here is carrots-only — no senior loses a benefit.", SANS, 12.5, INK, True)]], line_spacing=1.25)
rect(s, MX, Inches(6.35), Inches(2.4), Pt(2.2), ACCENT)
text(s, MX, Inches(6.55), Inches(11), Inches(0.4),
     [[mono("Prepared by Jeff Pinto · rvc-taxes.jeffpinto.com · every figure independently fact-checked", 8.5)]])

# 02 — WHAT JUST HAPPENED (timeline)
s = slide(); masthead(s, "The Squeeze")
text(s, MX, Inches(0.88), Inches(11.5), Inches(0.55), [[R("What just happened in Rockville Centre's schools", SERIF, 27, INK, True)]])
text(s, MX, Inches(1.48), Inches(11.6), Inches(0.55),
     [[R("An A-rated district — South Side High ranks #611 nationally; graduation 98–99% vs New York's 87% — chose between classrooms and solvency.", SERIF, 13.5, BODY, False, True)]], line_spacing=1.2)
TL = [("JAN 29, 2026", "District projects a $3.78M deficit for 2026-27: health insurance +10% (~$2M), out-of-district special-ed placements projected at $6.38M, total revenue growth only ~1% (~$1.4M)."),
      ("MAR 2026", "Final state levy cap lands at ≈2.1%. Reserves are thin: unappropriated fund balance $4.8M — 3.5% of budget against the 4% statutory ceiling; appropriated fund balance falling $2.8M → $1.6M."),
      ("APR 16, 2026", "Board adopts a $141.32M budget (+1.03%) closing the gap chiefly by eliminating ≈22.2 teaching positions and 40 teaching-assistant positions."),
      ("MAY 19, 2026", "Voters approve it 1,915–1,195.")]
y = Inches(2.25)
rect(s, MX + Inches(0.05), y, Pt(2), Inches(3.6), RULE)
for d, body in TL:
    rect(s, MX, y + Inches(0.06), Inches(0.12), Inches(0.12), ACCENT)
    text(s, MX + Inches(0.32), y, Inches(2.0), Inches(0.3), [[mono(d, 9, ACCENT, True)]])
    text(s, MX + Inches(2.45), y - Inches(0.02), Inches(9.9), Inches(0.85), [[R(body, SANS, 11.5, BODY)]], line_spacing=1.15)
    y += Inches(0.92)
dark_band(s, MX, Inches(6.02), W - MX*2, Inches(0.82))
text(s, MX + Inches(0.25), Inches(6.14), Inches(11.6), Inches(0.6),
     [[mono("Structural, not cyclical   ", 8.5, ACCENT_LT, True),
       R("Costs compound faster than the ~2% levy cap + ~15% state aid. A district can cut its way through one year, not a decade. The 2026-27 cuts are what “do nothing about housing” looks like.", SANS, 11.5, BG)]], line_spacing=1.15)
footrule(s, "Sources: LI Herald Jan–May 2026 · NYSED · OSC — full citations in appendix")

# 03 — THE SQUEEZE IN FOUR NUMBERS + BURNING PLATFORM
s = slide(); masthead(s, "The Squeeze")
text(s, MX, Inches(0.88), Inches(11.5), Inches(0.55), [[R("The squeeze, in four numbers", SERIF, 27, INK, True)]])
cw = (W - MX*2 - Inches(0.9)) / 4
cells = [("01 · Positions cut", "−62", "", "Teaching + TA positions eliminated to balance the 2026-27 budget", "LI Herald · Apr 2026"),
         ("02 · Enrollment", "−7", "%", "K-12 since 2015-16 (3,533 → 3,276; 3,260 projected for 2026-27)", "NYSED"),
         ("03 · Revenue", "+1", "%", "Total revenue growth — against +$2M for health insurance alone", "District · Mar 2026"),
         ("04 · Reserves", "3.5", "%", "Unappropriated fund balance; the statutory ceiling is 4%", "Asst. Supt. · Feb 2026")]
for i, (ey, k, u, c, src) in enumerate(cells):
    kpi_cell(s, MX + i*(cw + Inches(0.3)), Inches(1.75), cw, ey, k, u, c, src, kpi_size=34)
text(s, MX, Inches(4.12), Inches(11.8), Inches(0.75),
     [[R("Quality is not the question — NYSED's official per-pupil expenditure is $37,973 (2023-24), near the top of the peer set, and the results show it. ", SANS, 12, BODY),
       R("The question is whether a shrinking, aging base of families can keep voting for it.", SANS, 12, INK, True)]], line_spacing=1.25)
dark_band(s, MX, Inches(5.05), W - MX*2, Inches(1.7))
text(s, MX + Inches(0.25), Inches(5.22), Inches(11.7), Inches(0.35), [[mono("The burning platform", 9.5, ACCENT_LT, True)]])
text(s, MX + Inches(0.25), Inches(5.58), Inches(11.7), Inches(1.05),
     [[R("Nassau's school districts were built and staffed for more young families per capita than the population has aged into.", SERIF, 17, BG, True)],
      [R("Under a 2% levy cap, that mismatch balances only three ways: override votes, recurring cuts, or a bigger physical tax base. There is no fourth option.", SANS, 11.5, CREAM_DIM)]], line_spacing=1.2)
footrule(s, "NYSED financial transparency 2023-24 · LI Herald vote results")

# 04 — NOT JUST RVC + IF BENCHMARKS ARE MISSED
s = slide(); masthead(s, "The Squeeze")
text(s, MX, Inches(0.88), Inches(11.5), Inches(0.55), [[R("Rockville Centre is the rule, not the exception", SERIF, 27, INK, True)]])
kpi_cell(s, MX, Inches(1.70), Inches(3.7), "Nassau BOCES", "61", "%", "Of Nassau's 56 school districts lost enrollment over the past decade", "LI Herald · Apr 2025", kpi_size=40)
kpi_cell(s, MX + Inches(4.0), Inches(1.70), Inches(3.7), "W. Suffolk BOCES", "−6.3", "%", "Long Island public-school enrollment since 2017 (Nassau −4.0%)", "Bi-county study · 2024", kpi_size=40)
kpi_cell(s, MX + Inches(8.0), Inches(1.70), Inches(4.05), "Next door", "−16.5", "%", "Oyster Bay–East Norwich enrollment since 2015-16; Hewlett-Woodmere cut 11 teaching positions for 2025-26", "NYSED · LI Herald", kpi_size=40)
text(s, MX, Inches(4.0), Inches(11), Inches(0.3), [eyebrow("If the benchmarks are missed, the sequence is already written")])
seq = [("1 · OVERRIDES OR CUTS", "A 60% supermajority override on some of America's highest property taxes — or annual position cuts. Already happening: −62 (RVC), −11 (Hewlett-Woodmere)."),
       ("2 · EROSION & CLOSURES", "Larger classes, fewer electives, consolidations — then buildings. Northport closed two elementaries in 2021 (−10.6% enrollment); Long Beach parents fought one off in Jan 2025; Locust Valley's 2026-27 budget failed outright (revote June 16)."),
       ("3 · HOME-VALUE RISK", "School quality capitalizes into prices (Black 1999, QJE, and a large successor literature) — the premium that justifies RVC's taxes erodes with the schools.")]
y = Inches(4.4)
sw = (W - MX*2 - Inches(0.6)) / 3
for i, (t, d) in enumerate(seq):
    x = MX + i * (sw + Inches(0.3))
    rect(s, x, y, sw, Pt(2), RULE)
    text(s, x, y + Inches(0.10), sw, Inches(0.25), [[mono(t, 8.5, ACCENT, True)]])
    text(s, x, y + Inches(0.40), sw, Inches(1.6), [[R(d, SANS, 10.5, BODY)]], line_spacing=1.18)
footrule(s, "A county-wide coalition exists for the asking — slide 15")

# 05 — THE CAUSE
s = slide(); masthead(s, "The Cause")
text(s, MX, Inches(0.88), Inches(12), Inches(0.55), [[R("The village aged in place — the homes stopped turning over", SERIF, 26, INK, True)]])
cells = [("01 · Senior-owned", "27.9", "%", "Owner households headed by someone 65+ — 2,076 of 7,453", "ACS B25007"),
         ("02 · Senior share", "20.1", "%", "Of all RVC residents are 65 or older", "ACS S0101"),
         ("03 · Home value", "$819", "K", "Average RVC home, June 2026 (ACS 2020-24: $818.7K)", "Zillow · ACS B25077"),
         ("04 · Senior income", "~$80", "K", "Median 65+ household income, vs $151.9K village-wide", "ACS B19049")]
for i, (ey, k, u, c, src) in enumerate(cells):
    kpi_cell(s, MX + i*(cw + Inches(0.3)), Inches(1.70), cw, ey, k, u, c, src, kpi_size=34)
text(s, MX, Inches(4.05), Inches(11.8), Inches(0.8),
     [[R("Two thousand of the village's ~7,450 owner-occupied homes — the three- and four-bedroom houses that built the school system — are held by households past child-rearing age. Nationally, empty-nest boomers own ", SANS, 12, BODY),
       R("28% of large homes; millennials raising children own 14–16%", SANS, 12, INK, True), R(" (Redfin 2024/2026).", SANS, 12, BODY)]], line_spacing=1.25)
text(s, MX, Inches(5.05), Inches(11), Inches(0.3), [eyebrow("Nassau households including someone 65+ (ACS B11007)")])
bars = [("2014-18", 36.1), ("2019-23", 38.9), ("2020-24", 39.6)]
for i, (lbl, v) in enumerate(bars):
    yy = Inches(5.42) + Inches(0.4)*i
    text(s, MX, yy, Inches(1.1), Inches(0.3), [[mono(lbl, 8)]])
    rect(s, MX + Inches(1.25), yy + Inches(0.02), Inches(9.2), Inches(0.24), PAPER, line=RULESOFT)
    bw = Inches(9.2) * (v/45.0)
    rect(s, MX + Inches(1.25), yy + Inches(0.02), bw, Inches(0.24), ACCENT if i<2 else INK)
    text(s, MX + Inches(1.35) + bw, yy + Inches(0.02), Inches(1.0), Inches(0.25), [[mono(f"{v}%", 8.5, INK, True)]])
text(s, MX, Inches(6.62), Inches(11.8), Inches(0.3),
     [[mono("Households with own children under 18 over the same span: 35.5% → 34.7%", 8)]])
footrule(s, "ACS = American Community Survey 5-year estimates, Rockville Centre village & Nassau County")

# 06 — INVENTORY (dark quote — the single Locust Valley spot)
s = slide(dark=True)
text(s, MX, Inches(0.5), Inches(11), Inches(0.3), [[mono("The pipeline is inventory-starved, not demand-starved", 10, ACCENT_LT, True)]])
rect(s, MX, Inches(0.95), Inches(2.4), Pt(2), ACCENT_LT)
text(s, MX, Inches(1.7), Inches(1.2), Inches(2.0), [[R("“", SERIF, 110, ACCENT_LT, True)]])
text(s, MX + Inches(1.3), Inches(1.9), Inches(10.5), Inches(2.4),
     [[R("Five houses under a million dollars. In the whole district.", SERIF, 38, BG, True)]], line_spacing=1.05)
text(s, MX + Inches(1.3), Inches(4.0), Inches(10.3), Inches(1.6),
     [[R("Supt. Kristen Turnow, Locust Valley — “this is a Long Island problem, not a Locust Valley problem” — explaining a projected −15% enrollment slide (2,034 → 1,725 by 2027-28). Three weeks ago hers became the only Nassau district whose 2026-27 budget failed (850–765; revote June 16). Young families bid on the few large homes that list; the rest never list. Rockville Centre's version: 2,076 senior-owned homes, a frozen assessment roll, and a federal tax code that charges six figures for selling.", SANS, 13, CREAM_DIM)]], line_spacing=1.3)
text(s, MX + Inches(1.3), Inches(6.25), Inches(10), Inches(0.3), [[mono("LI Herald · Apr 2025 & May 2026 · NYSED · ACS B25007", 8, MONO_DIM)]])

# 07 — VOICES (the library, surfaced)
s = slide(); masthead(s, "The Voices")
text(s, MX, Inches(0.88), Inches(11.5), Inches(0.55), [[R("The rooms are telling you", SERIF, 27, INK, True)]])
voices = [("CHRIS ZUBLIONIS · SUPT., NORTH SHORE CSD · APR 2025",
           "You look at the students in senior year, and then you look at the incoming kindergarten, and it's like half. You can really see it going down."),
          ("LESLI DENINNO · PRESIDENT, RVC TEACHERS' ASSN. · MAR 2026",
           "The strength of this district is its teachers, staff and programs that support students. Cutting these is not the best practice. It is actively damaging to students and the community."),
          ("JEFF GREENFIELD · RVC RESIDENT, 56TH YEAR OF BUDGET HEARINGS · APR 2026",
           "It's probably one of the worst budget situations that I've seen in all those years. … You lost the confidence of the community as shepherds and stewards of our money.")]
y = Inches(1.7)
for who, q in voices:
    rect(s, MX, y, W - MX*2, Pt(1), RULESOFT)
    text(s, MX, y + Inches(0.14), Inches(0.4), Inches(0.6), [[R("“", SERIF, 30, ACCENT, True)]])
    text(s, MX + Inches(0.45), y + Inches(0.12), Inches(11.3), Inches(0.75),
         [[R(q, SERIF, 14.5, INK, False, True)]], line_spacing=1.22)
    text(s, MX + Inches(0.45), y + Inches(1.04), Inches(11.3), Inches(0.25), [[mono(who, 8, MUTED, True)]])
    y += Inches(1.46)
dark_band(s, MX, y + Inches(0.08), W - MX*2, Inches(0.62))
text(s, MX + Inches(0.25), y + Inches(0.20), Inches(11.6), Inches(0.4),
     [[mono("The library   ", 8.5, ACCENT_LT, True),
       R("30+ verbatim, sourced quotes — parents, teachers, superintendents, trustees — at rvc-taxes.jeffpinto.com/voices.html. Built for reuse: testimony, letters, op-eds.", SANS, 11.5, BG)]])
footrule(s, "All quotes verbatim from linked sources — LI Herald · News 12 · Patch · Northport Journal")

# 07 — FOUR LOCKS (2x2)
s = slide(); masthead(s, "The Lock-In")
text(s, MX, Inches(0.88), Inches(11.5), Inches(0.55), [[R("Four locks on the front door", SERIF, 27, INK, True)]])
locks = [("LOCK 01 · FEDERAL · CAP GAINS", "§121 frozen since 1997",
          "A 1990s buyer selling at ~$900K clears the $500K joint exclusion; a widowed senior filing single hits $250K twice as fast. AEI: ~1.9M senior-owned homes locked, $620B in gains.", "Fix: H.R. 1340 / S. 3332 — double & index"),
         ("LOCK 02 · STATE · TIMING TRAP", "Move wrong, lose a year's exemption",
          "Close at the wrong point in the assessment calendar and the senior exemption lapses for the year — a real, documented penalty for moving.", "Fix: S3309/A5288 — out of Aging 7–0"),
         ("LOCK 03 · STATE · CASH FLOW", "Fixed incomes hoard equity",
          "Seniors guard cash against carrying costs. Four states fix this with deferral: state pays, lien accrues at 5–8%, settled in full at sale (OR · MA · WA · TX).", "Fix: NY deferral enabling act"),
         ("LOCK 04 · LOCAL · INVENTORY", "Nothing to move into",
          "RVC is overwhelmingly single-family; the condo/ADU rungs of the ladder are missing, so selling means leaving the village after 40 years.", "Fix: ADUs, right-size zoning near the LIRR")]
gw = (W - MX*2 - Inches(0.5)) / 2
for i, (ey, h3, body, fix) in enumerate(locks):
    x = MX + (i % 2) * (gw + Inches(0.5)); y = Inches(1.7) + (i // 2) * Inches(2.45)
    rect(s, x, y, gw, Pt(2), RULE)
    text(s, x, y + Inches(0.10), gw, Inches(0.25), [[mono(ey, 8, ACCENT, True)]])
    text(s, x, y + Inches(0.38), gw, Inches(0.35), [[R(h3, SERIF, 17, INK, True)]])
    text(s, x, y + Inches(0.80), gw, Inches(1.0), [[R(body, SANS, 10.5, BODY)]], line_spacing=1.18)
    text(s, x, y + Inches(1.95), gw, Inches(0.3), [[mono(fix, 8, INK, True)]])
footrule(s, "§121 = 26 U.S.C. §121 · §467 = NY RPTL §467 income-tested senior exemption")

# 08 — MYTH-BUSTER (dark)
s = slide(dark=True)
text(s, MX, Inches(0.5), Inches(11), Inches(0.3), [[mono("Myth-buster for senior constituents", 10, ACCENT_LT, True)]])
rect(s, MX, Inches(0.95), Inches(2.4), Pt(2), ACCENT_LT)
text(s, MX, Inches(1.35), Inches(11.8), Inches(1.2), [[R("Moving within New York does not forfeit your benefits.", SERIF, 34, BG, True)]], line_spacing=1.05)
myths = [("ENHANCED STAR", "Follows you to any new NY primary residence — re-register, income-qualified, same credit."),
         ("§467 SENIOR EXEMPTION", "Applies at the new home wherever local jurisdictions adopted it — and Nassau waives its 12-month ownership wait for prior holders."),
         ("YOUR TAX BILL", "New York taxes current value, not purchase price (unlike California). Downsize to a smaller home and the bill falls with it — no Prop-13-style reset.")]
y = Inches(2.9)
for t, d in myths:
    rect(s, MX, y + Inches(0.05), Inches(0.12), Inches(0.12), ACCENT_LT)
    text(s, MX + Inches(0.35), y, Inches(2.9), Inches(0.5), [[mono(t, 9.5, ACCENT_LT, True)]])
    text(s, MX + Inches(3.4), y - Inches(0.03), Inches(9.1), Inches(0.85), [[R(d, SANS, 13, BG)]], line_spacing=1.2)
    y += Inches(1.02)
text(s, MX, Inches(6.35), Inches(11.8), Inches(0.6),
     [[R("The one real timing penalty — losing the exemption in the year of a move — is exactly what S3309/A5288 repairs.", SERIF, 14.5, CREAM_DIM, False, True)]])

# 09 — MORE 'STAY' INCENTIVES COMING (replaces "the lock is tightening")
s = slide(); masthead(s, "The Lock-In")
text(s, MX, Inches(0.88), Inches(12), Inches(0.55), [[R("More “stay” incentives are coming — with nothing for “go”", SERIF, 25, INK, True)]])
kpi_cell(s, MX, Inches(1.66), Inches(3.7), "Ch. 581 of 2025", "65", "%", "New maximum §467 exemption tier (income ≤ $47,000) — signed Dec 5, 2025; Nassau already applies it", "nysenate.gov · Nassau brochure", kpi_size=40)
kpi_cell(s, MX + Inches(4.0), Inches(1.66), Inches(3.7), "Separate 2025 law", "$75", "K", "§467 income ceiling from July 1, 2027 — a major expansion of the eligible pool, if localities opt in", "DTF 2025 legislative summary", kpi_size=40)
kpi_cell(s, MX + Inches(8.0), Inches(1.66), Inches(4.05), "NBER w25468", "100", "yrs", "Of Georgia exemptions, quadruple-difference design: age-based relief measurably increases seniors staying put", "Banzhaf, Mickey & Patrick · JUE 2021", kpi_size=40)
text(s, MX, Inches(3.95), Inches(11), Inches(0.3), [eyebrow("Each unpaired expansion makes three things harder")])
trio = [("FOR YOUNG FAMILIES", "Even fewer listings — relief raises the financial pull of staying, so fewer of the 2,076 senior-owned homes ever reach the market."),
        ("FOR CURRENT TAXPAYERS", "A larger §467 pool shifts a larger share of the fixed levy onto everyone else — working families most of all."),
        ("FOR SCHOOL DISTRICTS", "Slower turnover means faster enrollment decline — and the cuts-or-override arithmetic of slide 4 arrives sooner.")]
y = Inches(4.35)
for i, (t, d) in enumerate(trio):
    x = MX + i * (sw + Inches(0.3))
    rect(s, x, y, sw, Pt(2), RULE)
    text(s, x, y + Inches(0.10), sw, Inches(0.25), [[mono(t, 8.5, ACCENT, True)]])
    text(s, x, y + Inches(0.40), sw, Inches(1.4), [[R(d, SANS, 10.5, BODY)]], line_spacing=1.18)
text(s, MX, Inches(6.25), Inches(11.8), Inches(0.6),
     [[R("The ask is not to oppose relief — it is to pair every expansion with a fiscal note and a mobility counterpart.", SERIF, 15, BODY, False, True)]])
footrule(s, "Every §467 expansion decision is now pending locally — the data ask (slide 15) comes first")

# 10 — FOLLOW ONE HOME
s = slide(); masthead(s, "The Math, Done Honestly")
text(s, MX, Inches(0.88), Inches(11), Inches(0.55), [[R("Follow one home.", SERIF, 28, INK, True)]])
text(s, MX, Inches(1.46), Inches(11.8), Inches(0.5),
     [[R("A senior household sells; a young family buys. Four things happen — to four different balance sheets. None of them is “the district recaptures revenue.”", SERIF, 13.5, BODY, False, True)]], line_spacing=1.2)
oc = [("01 · Every year after the sale", "NY STATE SAVES", C_STATE, T_STATE, "$2,058", "/yr",
       "Enhanced STAR credit ($3,147) becomes Basic ($1,089). STAR is state-funded — the saver is Albany, forever.", "DTF 2025 final · RVC Class 1 · RPTL §1306-a"),
      ("02 · Once, at closing", "NY STATE COLLECTS", C_STATE, T_STATE, "~$3,300", "",
       "Transfer tax: 0.4% of an ~$820K sale, seller-paid (+1% buyer mansion tax at $1M+).", "Tax Law Art. 31 · ACS/Zillow"),
      ("03 · Every year after the sale", "NEIGHBORS RELIEVED", C_TAX, T_TAX, "≈$7,300", "/yr",
       "School tax un-shifted from neighbors at §467's 50% tier, on a verified ~$14.7K school bill. It's a sliding scale by income — 65% (≤$47K) down to 5% ($58.4K) — so ≈$730–$9,500 per home. The levy is fixed; exemptions move it around.", "OSC · Nassau brochure · Nassau LRV parcel tax table 2025-26"),
      ("04 · As the family invests", "DISTRICT GAINS", C_DIST, T_DIST, "+1.4", "students",
       "Per home (2 children × 70% school-age — stated assumption) into seats emptied by a ~7% enrollment decline; renovations grow the tax base, lawfully raising the levy cap.", "DTF growth factor · NYSED")]
gw2 = (W - MX*2 - Inches(0.5)) / 2
for i, (ey, tag, tc, tt, k, u, c, src) in enumerate(oc):
    x = MX + (i % 2) * (gw2 + Inches(0.5)); y = Inches(2.15) + (i // 2) * Inches(2.30)
    kpi_cell(s, x, y, gw2, ey, k, u, c, src, kpi_size=32, tag=tag, tagc=tc, tagt=tt)
footrule(s, "Full prose write-up: rvc-taxes.jeffpinto.com/fiscal-math.html")

# 11 — PER 100 + what it's not
s = slide(); masthead(s, "The Math, Done Honestly")
text(s, MX, Inches(0.88), Inches(11), Inches(0.55), [[R("Scale it: per 100 transitions", SERIF, 27, INK, True)]])
dark_band(s, MX, Inches(1.7), W - MX*2, Inches(1.95))
pw = (W - MX*2 - Inches(0.8)) / 4
pulls = [("$205.8K", "/yr", "STAR savings to NY State — recurring; funds every carrot in this deck"),
         ("~$330K", "", "One-time transfer-tax receipts to NY State"),
         ("n = ?", "", "§467 parcels relieved — the county does not publish this. That is Ask 3."),
         ("~140", "", "Students into existing seats at low marginal cost — while spare capacity lasts")]
for i, (v, u, l) in enumerate(pulls):
    x = MX + Inches(0.35) + i * pw
    runs = [R(v, SERIF, 26, BG, True)]
    if u: runs.append(R(u, SERIF, 14, ACCENT_LT, False, True))
    text(s, x, Inches(2.0), pw - Inches(0.25), Inches(0.5), [runs])
    text(s, x, Inches(2.55), pw - Inches(0.35), Inches(0.95), [[mono(l, 7.6, RULESOFT)]], line_spacing=1.25)
rect(s, MX, Inches(4.0), W - MX*2, Pt(2), RULE)
text(s, MX, Inches(4.15), Inches(11.6), Inches(0.3), [eyebrow("What this is not")])
text(s, MX, Inches(4.5), Inches(11.8), Inches(1.7),
     [[R("It is not district revenue. ", SANS, 12, INK, True),
       R("NY school districts levy a fixed dollar amount; exemptions shift it, they don't shrink it — “Every exemption granted on a property shifts the tax burden to the non-exempt properties” (NYS Comptroller) — and STAR is reimbursed by the State (RPTL §1306-a). Claims that exemptions “cost the schools millions” or that turnover “recaptures revenue” don't survive contact with the levy mechanics.", SANS, 12, BODY)],
      [R("", SANS, 6, BODY)],
      [R("The district's true gains: growth-factor headroom, students into empty seats, and voters with children in the system.", SERIF, 13.5, BODY, False, True)]], line_spacing=1.25)
footrule(s, "Capacity caveat: enrollment decline ≠ measured building capacity; assumptions labeled")

# 12 — BREAK-EVEN (new)
s = slide(); masthead(s, "The Math, Done Honestly")
text(s, MX, Inches(0.88), Inches(11.8), Inches(0.55), [[R("What would “no more cuts” actually take?", SERIF, 27, INK, True)]])
text(s, MX, Inches(1.46), Inches(11.8), Inches(0.5),
     [[R("Hold service constant and the arithmetic is one line: revenue growth must match cost growth. Sample calculation — illustrative, every input adjustable.", SERIF, 13, BODY, False, True)]], line_spacing=1.2)
dark_band(s, MX, Inches(2.05), W - MX*2, Inches(0.52))
text(s, MX + Inches(0.25), Inches(2.18), Inches(11.7), Inches(0.3),
     [[mono("cost growth  ≤  78% levy × (2% cap + physical base growth)  +  15% aid × aid growth  +  7% other", 9.5, ACCENT_LT, True)]])
bw_cells = [("At 3% cost growth", "~$1.8M", "/yr", "Structural gap with no base growth — revenue rises only ~1.7% (0.78×2% + 0.15×1%). The last two budgets balanced it with cuts (+1.63%, +1.03%).", "Identity · validated shares"),
            ("To close it in-cap", "~1.65", "%/yr", "Required PHYSICAL tax-base growth (renovations + new units — the DTF quantity change). On an assumed ~$7B base: ~$115M/yr of construction.", "Assumption-labeled: full value pending DTF"),
            ("100 transitions deliver", "~$15M", "", "Of quantity change (at ~$150K renovation each) — ≈13% of the in-cap requirement. Turnover is necessary, not sufficient.", "Stated assumptions"),
            ("So the dial set is", "4", "dials", "Turnover + new units (ADUs, condos) + enrollment-linked aid + the cost line. No single dial closes it — which is the case for pulling all four.", "Multi-variable by construction")]
for i, (ey, k, u, c, src) in enumerate(bw_cells):
    kpi_cell(s, MX + i*(cw + Inches(0.3)), Inches(2.85), cw, ey, k, u, c, src, kpi_size=26)
text(s, MX, Inches(5.35), Inches(11.8), Inches(1.3),
     [[R("Instrument it: ", SANS, 12, INK, True),
       R("rvc-taxes.jeffpinto.com/breakeven.html — sliders for every assumption (cost growth, cap, aid, turnover, renovation value, new units, district full value). Built to be recalibrated the day the county releases exemption counts and DTF full values — the same data ask that powers everything else in this deck.", SANS, 12, BODY)]], line_spacing=1.25)
footrule(s, "Illustrative — labeled assumptions; worked example in docs/BREAKEVEN_SAMPLE.md")

# 13 — ASKS: FEDERAL & STATE
s = slide(); masthead(s, "The Asks")
text(s, MX, Inches(0.88), Inches(11.5), Inches(0.5), [[R("The asks — federal & state", SERIF, 27, INK, True)]])
rows = [("FEDERAL", "H.R. 1340 / S. 3332 · IN COMMITTEE · 117+ COSPONSORS", "Cosponsor the More Homes on the Market Act",
         "Rep. Laura Gillen (D, NY-4) — lives in RVC, led the SALT fight, not yet on the bill. Sens. Gillibrand (Approps THUD RM; Aging RM) & Schumer (Minority Leader) for the Senate companion. Double & index §121."),
        ("STATE", "S3309 / A5288 · AGING 7–0 · APR 2026", "Keep a senior's exemption in the year of a move",
         "Sen. Patricia Canzoneri-Fitzpatrick (R, SD-9) — an RVC resident — to champion; Asm. Judy Griffin (D, AD-21) carries the companion. Add S3574 renewal-notice protection."),
        ("STATE", "NEW BILL · MODEL: OR · MA · WA · TX", "Senior school-tax deferral enabling act",
         "District opt-in; lien accrues at modest interest; taxes paid in full at sale. The cash-flow lock disappears at zero long-run cost to the levy. Draft with OSC input.")]
y = Inches(1.7)
for lab, status, h3, body in rows:
    rect(s, MX, y, W - MX*2, Pt(2 if lab=="FEDERAL" else 1), RULE if lab=="FEDERAL" else RULESOFT)
    text(s, MX, y + Inches(0.12), Inches(2.5), Inches(0.6), [[mono(lab, 9, ACCENT, True)], [mono(status, 7, MUTED)]], line_spacing=1.3)
    text(s, MX + Inches(2.75), y + Inches(0.10), Inches(9.4), Inches(0.35), [[R(h3, SERIF, 17, INK, True)]])
    text(s, MX + Inches(2.75), y + Inches(0.48), Inches(9.4), Inches(0.85), [[R(body, SANS, 11, BODY)]], line_spacing=1.18)
    y += Inches(1.52)
dark_band(s, MX, y + Inches(0.06), W - MX*2, Inches(0.6))
text(s, MX + Inches(0.25), y + Inches(0.17), Inches(11.6), Inches(0.4),
     [[mono("Framing that works   ", 8.5, ACCENT_LT, True),
       R("Senior mobility + young-family affordability in one bill — and Albany funds the carrots out of its own STAR savings.", SANS, 11.5, BG)]])
footrule(s, "Officeholders & statuses verified June 2026 — congress.gov · nysenate.gov")

# 14 — ASKS: COUNTY & LOCAL
s = slide(); masthead(s, "The Asks")
text(s, MX, Inches(0.88), Inches(11.5), Inches(0.5), [[R("The asks — county & local: data before dollars", SERIF, 27, INK, True)]])
rows = [("COUNTY", "BEFORE THE 65% TIER & THE 2027 $75K CAP", "Publish §467 / Enhanced STAR counts & exempted value, by school district, annually",
         "County Executive Bruce Blakeman, Nassau Legislature, Dept. of Assessment. Require a fiscal note before any exemption expansion; fund a means-tested “Silver-to-Gold” moving-grant pilot ($5–10K)."),
        ("VILLAGE + BOE", "ROCKVILLE CENTRE", "Confirm the district's own §467 posture; build the downsizing ladder",
         "Mayor Francis X. Murray; RVC Board of Education (incl. Trustee-elect Mary Beth Joyce, seated July 2026): publish the adopting resolution and current tiers; pilot moving assistance ($25–50K/yr); pursue ADU / right-size supply near the LIRR.")]
y = Inches(1.7)
first = True
for lab, status, h3, body in rows:
    rect(s, MX, y, W - MX*2, Pt(2 if first else 1), RULE if first else RULESOFT); first = False
    text(s, MX, y + Inches(0.12), Inches(2.5), Inches(0.7), [[mono(lab, 9, ACCENT, True)], [mono(status, 7, MUTED)]], line_spacing=1.3)
    text(s, MX + Inches(2.75), y + Inches(0.10), Inches(9.4), Inches(0.6), [[R(h3, SERIF, 16, INK, True)]], line_spacing=1.05)
    text(s, MX + Inches(2.75), y + Inches(0.72), Inches(9.4), Inches(0.85), [[R(body, SANS, 11, BODY)]], line_spacing=1.18)
    y += Inches(1.75)
text(s, MX, y + Inches(0.05), Inches(11), Inches(0.3), [eyebrow("Twelve-month sequencing")])
seq2 = [("NOW–SEPT 2026", "County data publication · district §467 confirmation · coalition letters from the 61%-club"),
        ("FALL 2026", "Deferral bill drafted with OSC · H.R. 1340 cosponsor push pre-election"),
        ("2027 SESSION", "S3309/A5288 + deferral act pass before the July 2027 $75K-cap decision — relief paired with mobility")]
for i, (d, t) in enumerate(seq2):
    x = MX + i * (sw + Inches(0.3))
    rect(s, x, y + Inches(0.4), sw, Pt(2), RULE)
    text(s, x, y + Inches(0.5), sw, Inches(0.25), [[mono(d, 8.5, ACCENT, True)]])
    text(s, x, y + Inches(0.78), sw, Inches(0.8), [[R(t, SANS, 10, BODY)]], line_spacing=1.15)
footrule(s, "Coalition-ready: Hewlett-Woodmere · Locust Valley · OBEN · Long Beach · the rest of the 61%")

# 15 — PRECEDENTS
s = slide(); masthead(s, "Precedents & Evidence")
text(s, MX, Inches(0.88), Inches(12), Inches(0.5), [[R("Assembled from things that already work", SERIF, 27, INK, True)]])
rows = [("DEFERRAL · OR 1963 · MA 41A · WA · TX", "Districts made whole; seniors cash-flow-safe", "State pays; a lien — not a subsidy — does the work, repaid with 5–8% interest at sale. Voluntary. Oregon's program is self-financing."),
        ("CALIFORNIA · PROP 19 · 2020", "Senior mobility wins at the ballot box", "Base-transfer for 55+, statewide, up to 3 moves — narrowly approved by voters; LAO scored the package revenue-positive for schools. The lesson for NY is the politics, not the mechanism: NY taxes current value, so our version is benefit-continuity (S3309), not basis-porting."),
        ("MINNEAPOLIS FED · 2024", "Property taxes reduce lock-in, shift ownership younger", "Recurring taxes capitalize into lower prices and raise turnover toward young families — NY's turnover dividend is unusually large once artificial locks come off."),
        ("NORTHPORT · 2021 · THREE VILLAGE · 2025", "The endgame is current, not historical", "Northport–East Northport closed two elementaries in Aug 2021 (enrollment −10.6% + LIPA step-down; ~$7M/yr saved); Three Village consolidated grades district-wide in 2025; Long Beach is Nassau's nearest miss (East Elementary saved by 4,000 signatures, Jan 2025). The cycle is now Rockville Centre's to manage, with better tools.")]
y = Inches(1.62)
for ey, h3, body in rows:
    rect(s, MX, y, W - MX*2, Pt(1), RULESOFT)
    text(s, MX, y + Inches(0.1), Inches(3.2), Inches(0.55), [[mono(ey, 7.8, ACCENT, True)]], line_spacing=1.25)
    text(s, MX + Inches(3.45), y + Inches(0.08), Inches(8.7), Inches(0.3), [[R(h3, SERIF, 14.5, INK, True)]])
    text(s, MX + Inches(3.45), y + Inches(0.4), Inches(8.7), Inches(0.7), [[R(body, SANS, 10, BODY)]], line_spacing=1.15)
    y += Inches(1.13)
text(s, MX, y + Inches(0.08), Inches(11.8), Inches(0.7),
     [[R("What we are not proposing: ", SANS, 11.5, INK, True),
       R("no benefit cuts, no forced sales, no “empty-nester tax.” Opt-in for the senior; revenue-positive for the State.", SANS, 11.5, BODY)]], line_spacing=1.2)
footrule(s, "NBER w25468 evidence on exemption retention effects — slide 10")

# 16 — THE UNASKED QUESTION: DISTRICT LINES (new)
s = slide(); masthead(s, "The Horizon Lever")
text(s, MX, Inches(0.88), Inches(12), Inches(0.55), [[R("The unasked question: the map itself", SERIF, 27, INK, True)]])
text(s, MX, Inches(1.48), Inches(11.8), Inches(0.55),
     [[R("Why does one county run 56 school districts whose lines match no town, no village — and never move?", SERIF, 13.5, BODY, False, True)]], line_spacing=1.2)
maprows = [("THE LINES ARE ACCIDENTS", "Nassau's district boundaries descend from 19th-century common-school districts. They are not coterminous with towns or villages (RVC UFSD ≠ Village of RVC), and they have been treated as carved in stone for a century."),
           ("ALBANY ALREADY PAYS FOR THIS", "NYSED's reorganization framework offers merging districts a multi-year incentive operating-aid bonus plus enhanced building aid — money on the table that Nassau districts have left untouched (NYSED, Guide to the Reorganization of School Districts)."),
           ("SHORT OF MERGER", "Utilization can rebalance across lines without redrawing them: cross-district tuitioning, shared programs through BOCES, and boundary-alteration petitions for edge neighborhoods."),
           ("WHY IT BELONGS ON THE TABLE", "A county where 61% of districts are shrinking should not treat the map as scripture. If the families can't come to the seats, the seats question eventually comes to the map — better to study it with data than meet it in a crisis.")]
y = Inches(2.15)
for t, d in maprows:
    rect(s, MX, y, W - MX*2, Pt(1), RULESOFT)
    text(s, MX, y + Inches(0.09), Inches(3.2), Inches(0.5), [[mono(t, 8, ACCENT, True)]], line_spacing=1.25)
    text(s, MX + Inches(3.45), y + Inches(0.07), Inches(8.7), Inches(0.85), [[R(d, SANS, 10.5, BODY)]], line_spacing=1.16)
    y += Inches(1.02)
dark_band(s, MX, y + Inches(0.1), W - MX*2, Inches(0.6))
text(s, MX + Inches(0.25), y + Inches(0.21), Inches(11.6), Inches(0.4),
     [[mono("Candor   ", 8.5, ACCENT_LT, True),
       R("This is the political hot potato of the set — which is why the ask here is only the study and the data, not the redraw.", SANS, 11.5, BG)]])
footrule(s, "Same county data ask (slide 15) powers this conversation — utilization + exemptions by district")

# 17 — CLOSE (dark CTA)
s = slide(dark=True)
text(s, MX, Inches(0.55), Inches(11), Inches(0.3), [[mono("The close", 10, ACCENT_LT, True)]])
rect(s, MX, Inches(1.0), Inches(2.4), Pt(2), ACCENT_LT)
text(s, MX, Inches(1.55), Inches(12.0), Inches(2.4),
     [[R("One county. One data ask.", SERIF, 46, BG, True)], [R("One Albany agenda.", SERIF, 46, BG, True)]], line_spacing=1.05)
text(s, MX, Inches(3.6), Inches(11.5), Inches(1.4),
     [[R("Each budget cycle under the cap starts ~$2–4M behind. The 2026-27 answer was 62 classroom positions. The 2027-28 question arrives with a higher exemption ceiling, fewer students, and an older electorate — every year of delay makes the carrots more expensive and the politics harder.", SERIF, 16, CREAM_DIM, False, True)]], line_spacing=1.3)
kv = [("THE DECK", "rvc-taxes.jeffpinto.com"), ("THE MECHANICS", "/fiscal-math.html"), ("THE INSTRUMENT", "/breakeven.html"), ("THE RECEIPTS", "/validation.html")]
for i, (k, v) in enumerate(kv):
    x = MX + i * Inches(3.0)
    text(s, x, Inches(5.55), Inches(2.8), Inches(0.25), [[mono(k, 8, ACCENT_LT, True)]])
    text(s, x, Inches(5.85), Inches(2.9), Inches(0.3), [[R(v, SANS, 12, BG, True)]])
rect(s, MX, Inches(6.5), W - MX*2, Pt(1), RGBColor(0x3A, 0x3F, 0x48))
text(s, MX, Inches(6.62), Inches(11.8), Inches(0.4),
     [[mono("Prepared by Jeff Pinto · June 2026 · every figure independently fact-checked against primary sources", 8.5, MONO_DIM)]])

# =========================================================================
# APPENDIX — lean: validated facts + sources only
def appendix_head(s, title):
    masthead(s, "Appendix")
    text(s, MX, Inches(0.85), Inches(12), Inches(0.5), [[R(title, SERIF, 23, INK, True)]])

def table(s, x, y, w, col_widths, header, rows, fs=9.0, hfs=8.0, row_h=0.34):
    n_rows, n_cols = len(rows) + 1, len(header)
    gt = s.shapes.add_table(n_rows, n_cols, x, y, w, Inches(row_h * n_rows)).table
    for j, cw_ in enumerate(col_widths): gt.columns[j].width = Inches(cw_)
    for j, htxt in enumerate(header):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = INK
        c.margin_left = c.margin_right = Inches(0.06); c.margin_top = c.margin_bottom = Inches(0.02)
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = htxt.upper()
        r.font.name = MONO; r.font.size = Pt(hfs); r.font.bold = True; r.font.color.rgb = BG
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            c = gt.cell(i + 1, j); c.fill.solid()
            c.fill.fore_color.rgb = PAPER if i % 2 == 0 else BG
            c.margin_left = c.margin_right = Inches(0.06); c.margin_top = c.margin_bottom = Inches(0.02)
            p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = str(cell)
            r.font.name = SANS; r.font.size = Pt(fs); r.font.color.rgb = BODY
            c.text_frame.word_wrap = True
    return gt

# A1 — divider
s = slide()
text(s, MX, Inches(2.4), Inches(11), Inches(0.4), [eyebrow("Appendix · for reading, not presenting")])
text(s, MX, Inches(2.85), Inches(12), Inches(1.4), [[R("The numbers, and where they come from.", SERIF, 44, INK, True)]], line_spacing=1.0)
text(s, MX, Inches(4.35), Inches(11.4), Inches(1.8),
     [[R("Every figure in this deck was checked against primary sources on June 9, 2026: NYSED, NYS DTF, the State Comptroller, the Census Bureau (ACS), IRS SOI, congress.gov, nysenate.gov, Nassau County, and the named press coverage. The appendix carries the validated reference set — mechanics, tier tables, precedents, and citations. The full fact-check (per-claim verdicts, with everything that didn't survive) lives at rvc-taxes.jeffpinto.com/validation.html.", SANS, 12.5, BODY)]], line_spacing=1.3)
footrule(s, "Methodology: claims traced to statute, agency table, or named article; untraceable claims labeled, not repeated")

# A2 — levy mechanics
s = slide(); appendix_head(s, "Levy mechanics & STAR plumbing — the reference card")
mech = [("THE LEVY", "Districts collect a voter-approved, state-capped dollar amount, apportioned over taxable assessed value. Rate = levy ÷ base. Exemptions shrink the base; the rate rises; the district is whole. (OSC: “Every exemption granted on a property shifts the tax burden to the non-exempt properties.”)"),
        ("STAR", "RPTL §1306-a makes STAR savings “a state charge” — districts are reimbursed in full. Post-2015 buyers receive a state-paid credit check instead. Enhanced→Basic deltas accrue to Albany: $3,147.01 → $1,089 in RVC (Class 1, 2025 final)."),
        ("THE CAP", "Allowable levy = prior levy × tax-base growth factor (physical/quantity change only — construction and renovations, never price appreciation) × min(2%, CPI) + exclusions. Turnover-driven construction is the one lawful district-revenue channel."),
        ("ASSESSMENT", "Nassau Class 1 LOA = 0.1%; the roll is frozen for a sixth consecutive year (2026-27 tentative). Sales do NOT trigger reassessment — NY has no acquisition-value system (RPTL §305)."),
        ("TRANSFER TAXES", "NYS: 0.4% seller-paid ($2/$500) + 1% buyer mansion tax ≥$1M — both state revenue. Nassau has no county transfer tax (its 1999 authorization expired Jan 31, 2001).")]
y = Inches(1.55)
for k, v in mech:
    rect(s, MX, y, W - MX*2, Pt(1), RULESOFT)
    text(s, MX, y + Inches(0.08), Inches(2.2), Inches(0.4), [[mono(k, 9, ACCENT, True)]])
    text(s, MX + Inches(2.45), y + Inches(0.06), Inches(9.7), Inches(0.92), [[R(v, SANS, 9.8, BODY)]], line_spacing=1.16)
    y += Inches(1.04)
footrule(s, "OSC exemption brief · RPTL §1306-a · OSC tax-cap glossary · DTF transfer-tax pages")

# A3 — §467 tiers (validated table only)
s = slide(); appendix_head(s, "The §467 sliding scale — Nassau County (county · town · school)")
table(s, MX, Inches(1.5), Inches(6.3), [3.0, 3.3],
      ["2025 income", "Exemption"],
      [("≤ $47,000", "65%  (new — Ch. 581 of 2025)"), ("$47,001 – $47,999", "60%"), ("$48,000 – $48,999", "55%"),
       ("$49,000 – $50,000", "50%  (statutory base max)"), ("$50,001 – $52,999", "45% / 40% / 35%"),
       ("$53,000 – $55,699", "30% / 25% / 20%"), ("$55,700 – $57,499", "15% / 10%"),
       ("$57,500 – $58,399", "5%  (floor tier — the ceiling for ANY exemption)"), ("> $58,399", "0%")],
      fs=9.5, row_h=0.42)
text(s, MX + Inches(6.9), Inches(1.55), Inches(5.5), Inches(4.6),
     [[R("Who this reaches in RVC", SERIF, 14, INK, True)],
      [R("The income ceiling sits well below the village's 65+ median household income (~$80K, ACS B19049) — the eligible pool is real but a minority of the 2,076 senior-owned homes. The exact count, by tier, is the county data ask.", SANS, 10.5, BODY)],
      [R("", SANS, 6, BODY)],
      [R("Coming: $75,000 income cap", SERIF, 14, INK, True)],
      [R("Effective July 1, 2027 (separate 2025 legislation; local opt-in). Combined with the 65% tier, the eligible pool and the burden shift both grow — which is why the data ask precedes any local opt-in vote.", SANS, 10.5, BODY)],
      [R("", SANS, 6, BODY)],
      [R("Open item", SERIF, 14, INK, True)],
      [R("RVC UFSD's own adopting resolution and current ceiling — county materials say all Nassau districts participate; confirming with the district clerk (FOIL drafted).", SANS, 10.5, BODY)]],
     line_spacing=1.2)
footrule(s, "Nassau Dept. of Assessment brochure rev. 3-26 · DTF senior-exemption pages · DTF 2025 leg. summary")

# A4 — deferral
s = slide(); appendix_head(s, "Deferral precedents — the model for a NY enabling act")
table(s, MX, Inches(1.5), W - MX*2, [2.2, 2.4, 2.2, 5.3],
      ["Program", "Interest", "Repaid", "Mechanics"],
      [("Oregon (ORS 311.666–701, since 1963)", "6% simple", "Sale / death / move", "State DOR pays the county each Nov 15; lien recorded; income cap ~$70K indexed; recertify every 2 years; self-financing"),
       ("Massachusetts (Cl. 41A)", "8% (local option lower)", "Sale / death", "65+; local income limit; deferral + interest capped at 50% of fair cash value; 16% post-death until paid"),
       ("Washington (RCW 84.38)", "5%", "Sale / death", "61+/disabled; income threshold indexed to county median; deferral capped at 80% of equity"),
       ("Texas (Tax Code §33.06)", "5%", "Sale / estate", "65+/disabled; tax collection simply defers; widely used, simple affidavit"),
       ("Proposed: New York", "OSC-set", "At sale", "District opt-in for the school-tax share; lien accrues; the levy is made whole — the cash-flow lock (slide 8, Lock 03) disappears at zero long-run cost")],
      fs=9.0, row_h=0.52)
text(s, MX, Inches(4.85), Inches(11.8), Inches(0.8),
     [[R("Design note: deferral pairs naturally with S3309/A5288 — continuity for those who move, deferral for those who stay until they're ready. Both are senior-protective; neither costs the district its levy.", SERIF, 13, BODY, False, True)]], line_spacing=1.25)
footrule(s, "Oregon DOR · Mass. DLS guide · RCW 84.38 · Texas Comptroller")

# A5 — roster (facts only)
s = slide(); appendix_head(s, "June 2026 roster & live vehicles")
table(s, MX, Inches(1.5), W - MX*2, [2.6, 4.4, 5.1],
      ["Seat / vehicle", "Holder / status — June 2026", "Relevance"],
      [("NY-4 (US House)", "Rep. Laura Gillen (D) — RVC resident", "SALT-relief record; June 23, 2026 primary; not yet on H.R. 1340 — the ask"),
       ("US Senate", "Schumer (Minority Leader) · Gillibrand", "Gillibrand: Approps THUD Ranking Member; Special Committee on Aging RM"),
       ("NY SD-9", "Sen. Patricia Canzoneri-Fitzpatrick (R)", "RVC resident; seat on the ballot Nov 2026; natural S3309 champion"),
       ("NY AD-21", "Asm. Judy Griffin (D)", "District covers RVC, Lynbrook, Malverne, Baldwin; A5288 companion"),
       ("Nassau County Exec", "Bruce Blakeman (R)", "Re-elected Nov 2025; running for governor; owns the data ask"),
       ("RVC Village / BOE", "Mayor Francis X. Murray · Trustee-elect M.B. Joyce (July 2026)", "Village runs its own assessor & tax office; BOE owns the district §467 posture"),
       ("H.R. 1340 / S. 3332", "Ways & Means / Senate Finance", "Panetta-Kelly · Cornyn-Bennet-Schiff; double & index §121"),
       ("S3309 / A5288", "Reported from Aging 7–0 (Apr 2026)", "Exemption continuity in a move year (Palumbo)"),
       ("Ch. 581 of 2025", "Law — signed Dec 5, 2025", "65% §467 tier; income cap $75,000 eff. July 1, 2027 (separate 2025 law)")],
      fs=8.6, row_h=0.44)
footrule(s, "Verified against congress.gov · nysenate.gov · official sites")

# A6 — coalition (validated facts)
s = slide(); appendix_head(s, "The coalition map — validated district facts (NYSED 2024-25)")
table(s, MX, Inches(1.5), W - MX*2, [2.7, 2.0, 7.4],
      ["District", "K-12 (NYSED)", "Validated facts"],
      [("Hewlett-Woodmere", "2,650", "11 teaching positions cut for 2025-26; elementary enrollment projected −200 vs 2022-23"),
       ("Locust Valley", "1,754", "−15% projected by 2027-28 (2,034 → 1,725); “five homes under $1M” — Supt. Turnow; 2026-27 budget FAILED 850–765 (revote 6/16/26)"),
       ("Oyster Bay-E. Norwich", "1,310", "−16.5% since 2015-16; Supt. Ianni credits universal pre-K for stabilization"),
       ("Long Beach", "3,265", "Fought a year-long closure debate over East Elementary; community advocacy kept it open (2025)"),
       ("Levittown", "7,061", "−2% since 2022-23; $274.1M budget (2025-26)"),
       ("Lynbrook / Malverne / Massapequa / Great Neck", "2,725 / 1,801 / 6,460 / 6,600", "Budgets at or near their caps ($110M / $73.9M max-cap / $243M); Uniondale 65% Hispanic, 32% Black"),
       ("County-wide", "—", "61% of 56 districts declined over the decade (Nassau BOCES); LI −6.3% since 2017 (W. Suffolk BOCES)")],
      fs=8.6, row_h=0.52)
text(s, MX, Inches(5.6), Inches(11.8), Inches(0.6),
     [[R("One county, one data ask, one Albany agenda — the same brief works for every district above.", SERIF, 13, BODY, False, True)]])
footrule(s, "NYSED enrollment pages per district · LI Herald · W. Suffolk BOCES bi-county study")

# A7 — sources 1
s = slide(); appendix_head(s, "Sources (1 of 2) — district, demographics, housing")
src1 = [("RVC BUDGET PROCESS", "LI Herald: “$3.78M deficit” (Feb 6, 2026) · board-scrutiny (Feb 13) · “$141.3M adopted” (Apr 16) · “Budget passes” (May 19) · M.B. Joyce substack, “Understanding the 2026/2027 budget” (Mar 18)"),
        ("NYSED", "Enrollment (instid 800000049383): 3,533 → 3,276 · grad rate 98% · per-pupil expenditure $37,973 (2023-24) · district pages for all coalition districts"),
        ("CENSUS / ACS", "B25007 (2,076 senior-owner households) · B19049 (incomes by age) · B19013 ($144,516 → $151,938) · B25077 ($818,700) · B11007 (Nassau 38.9% → 39.6%) · IRS SOI TY2022 (itemization)"),
        ("HOUSING", "Zillow RVC home values (June 2026) · Redfin “Empty Nesters” (Jan 2024; 2026 refresh) · ATTOM property-tax reports (2023–2025)"),
        ("ENROLLMENT CONTEXT", "Nassau BOCES via LI Herald (61% of 56 districts) · W. Suffolk BOCES bi-county study (LI −6.3% since 2017) · Northport Journal/Patch (two 2021 closures) · LI Herald (Locust Valley budget failure, May 2026) · voices library: /voices.html")]
y = Inches(1.55)
for k, v in src1:
    rect(s, MX, y, W - MX*2, Pt(1), RULESOFT)
    text(s, MX, y + Inches(0.07), Inches(2.6), Inches(0.4), [[mono(k, 8.5, ACCENT, True)]])
    text(s, MX + Inches(2.85), y + Inches(0.06), Inches(9.3), Inches(0.85), [[R(v, SANS, 9.6, BODY)]], line_spacing=1.18)
    y += Inches(0.99)
footrule(s, "Clickable links: rvc-taxes.jeffpinto.com (deck p.8) and /validation.html")

# A8 — sources 2 + contact
s = slide(); appendix_head(s, "Sources (2 of 2) — law, mechanics, research · contact")
src2 = [("STATUTE & MECHANICS", "26 U.S.C. §121 · NY RPTL §467, §1306-a, §305 · Tax Law Art. 31 · OSC “Property Tax Exemptions in NYS” · OSC tax-cap glossary (growth factor) · Nassau §467 brochure (rev. 3-26) · DTF STAR tables (2025 final) & 2025 legislative summary ($75K cap)"),
        ("BILLS", "H.R. 1340 (Panetta/Kelly) · S. 3332 (Cornyn/Bennet/Schiff) · S3309/A5288 (Palumbo — Aging 7–0 Apr 2026) · S3574 · S5175A → Ch. 581 of 2025 (signed Dec 5, 2025)"),
        ("RESEARCH", "AEI Housing Center (Pinto): ~1.9M senior homes, $620B gains · Banzhaf, Mickey & Patrick, NBER w25468 / J. Urban Economics 2021 · Horwich, Minneapolis Fed (Nov 14, 2024) · Black (1999), QJE — school quality capitalization · CA BOE Prop 19 + LAO"),
        ("REORGANIZATION & DEFERRAL", "NYSED, Guide to the Reorganization of School Districts (incentive operating & building aid) · Oregon ORS 311.666–701 · Mass. G.L. c.59 §5 Cl. 41A · Washington RCW 84.38 · Texas Tax Code §33.06")]
y = Inches(1.55)
for k, v in src2:
    rect(s, MX, y, W - MX*2, Pt(1), RULESOFT)
    text(s, MX, y + Inches(0.07), Inches(2.6), Inches(0.5), [[mono(k, 8.5, ACCENT, True)]])
    text(s, MX + Inches(2.85), y + Inches(0.06), Inches(9.3), Inches(0.9), [[R(v, SANS, 9.6, BODY)]], line_spacing=1.18)
    y += Inches(1.06)
dark_band(s, MX, y + Inches(0.05), W - MX*2, Inches(0.85))
text(s, MX + Inches(0.25), y + Inches(0.16), Inches(11.6), Inches(0.6),
     [[mono("Contact & next steps   ", 8.5, ACCENT_LT, True),
       R("Jeff Pinto · rvc-taxes.jeffpinto.com — FOIL to Nassau Assessment (exemption counts by district) · district-clerk confirmation of RVC's §467 posture · coalition letters from the 61%-club.", SANS, 11, BG)]], line_spacing=1.2)
footrule(s, "Slides 1–18 are the live deck; this appendix is the validated reference set")

# ---- save ----------------------------------------------------------------
out = os.path.join(os.path.dirname(__file__), "..", "site", "RVC_Legislator_Deck.pptx")
prs.save(os.path.abspath(out))
print("wrote", os.path.abspath(out), "slides:", PAGE_N)
